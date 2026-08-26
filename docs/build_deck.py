from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx.oxml.ns import qn

# palette — matches the HTML deck
INK   = RGBColor(0x25,0x31,0x3A)
INK2  = RGBColor(0x36,0x45,0x4F)
MUTE  = RGBColor(0x6C,0x7B,0x88)
FAINT = RGBColor(0x9A,0xA8,0xB3)
LINE  = RGBColor(0xD3,0xD8,0xDD)
LINE2 = RGBColor(0xC2,0xCA,0xD1)
PANEL = RGBColor(0xEE,0xF1,0xF3)
STAGE = RGBColor(0xFB,0xFB,0xFC)
WHITE = RGBColor(0xFF,0xFF,0xFF)
AMBER = RGBColor(0xB2,0x6B,0x12)
AMBD  = RGBColor(0x8A,0x52,0x10)
AMBSOFT=RGBColor(0xF6,0xE9,0xD6)
AMBLINE=RGBColor(0xE3,0xB1,0x73)
CLEAN = RGBColor(0x4B,0x66,0x50); CLEANS=RGBColor(0xE7,0xED,0xE8)
DRIFT = AMBER;                    DRIFTS=AMBSOFT
ERROR = RGBColor(0x8C,0x3B,0x2E); ERRORS=RGBColor(0xF3,0xE2,0xDE)
SANS="DejaVu Sans"; MONO="DejaVu Sans Mono"

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
W,H=13.333,7.5; Mx=1.05; CW=W-2*Mx

def slide(bg=STAGE):
    s=prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb=bg
    # inner frame
    fr=s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.42),Inches(0.42),Inches(W-0.84),Inches(H-0.84))
    fr.fill.background(); fr.line.color.rgb=LINE; fr.line.width=Pt(1); fr.shadow.inherit=False
    return s

def box(s,x,y,w,h,fill=None,line=None,lw=1.0,dash=None,laccent=None):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb=line; sh.line.width=Pt(lw)
    sh.shadow.inherit=False
    if laccent:
        t=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(w),Pt(3))
        t.fill.solid(); t.fill.fore_color.rgb=laccent; t.line.fill.background(); t.shadow.inherit=False
    return sh

def txt(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,sp=4,leading=None):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=tf.margin_right=tf.margin_top=tf.margin_bottom=0
    first=True
    for line in runs:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        p.alignment=align; p.space_after=Pt(sp); p.space_before=Pt(0)
        if leading: p.line_spacing=leading
        for (t,sz,c,b,f) in line:
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=b
            r.font.name=f; r.font.color.rgb=c
    return tb

def kicker(s,text,dark=False):
    t=box(s,Mx,0.72,0.28,0.028,AMBER)
    txt(s,Mx+0.4,0.6,CW,0.4,[[(text.upper(),11,(FAINT if dark else MUTE),False,MONO)]])

def foot(s,label,n,dark=False):
    c=RGBColor(0x7C,0x8B,0x95) if dark else FAINT
    txt(s,Mx,H-0.86,5,0.3,[[(label.upper(),9.5,c,False,MONO)]])
    txt(s,W-Mx-2.4,H-0.86,2.4,0.3,[[(f"{n:02d}",9.5,(RGBColor(0xC4,0xCE,0xD5) if dark else INK2),True,MONO),(f" / 18",9.5,c,False,MONO)]],align=PP_ALIGN.RIGHT)

def title(s,parts,y=1.5,size=33,x=None,w=None):
    # parts: list of (text,color) segments on possibly one line
    x=Mx if x is None else x; w=CW if w is None else w
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(1.6)); tf=tb.text_frame
    tf.word_wrap=True; p=tf.paragraphs[0]; p.line_spacing=1.02
    for (t,c) in parts:
        r=p.add_run(); r.text=t; r.font.size=Pt(size); r.font.bold=True; r.font.name=SANS; r.font.color.rgb=c
    return tb

def bullets(s,x,y,w,items,size=15,gap=0.5,color=MUTE,dot=AMBER):
    for i,it in enumerate(items):
        box(s,x,y+i*gap+0.09,0.09,0.09,dot if (isinstance(it,tuple)) else MUTE)
        t=it[0] if isinstance(it,tuple) else it
        col=it[1] if isinstance(it,tuple) else color
        txt(s,x+0.32,y+i*gap,w-0.32,gap,[[(t,size,col,False,SANS)]])

def codebox(s,x,y,w,lines,size=13):
    h=0.34+len(lines)*(size+9)/72.0
    box(s,x,y,w,h,PANEL,None); box(s,x,y,0.045,h,AMBER)
    tb=s.shapes.add_textbox(Inches(x+0.25),Inches(y+0.16),Inches(w-0.45),Inches(h-0.3))
    tf=tb.text_frame; tf.word_wrap=True; first=True
    for segs in lines:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        p.space_after=Pt(2); p.line_spacing=1.15
        for (t,c,b) in segs:
            r=p.add_run(); r.text=t; r.font.size=Pt(size); r.font.name=MONO; r.font.color.rgb=c; r.font.bold=b
    return y+h

def arrow(s,x1,y1,x2,y2,color=MUTE,dash=False,w=1.4):
    cn=s.shapes.add_connector(2,Inches(x1),Inches(y1),Inches(x2),Inches(y2))
    cn.line.color.rgb=color; cn.line.width=Pt(w); cn.shadow.inherit=False
    le=cn.line._get_or_add_ln()
    he=le.makeelement(qn('a:tailEnd'),{'type':'triangle','w':'med','len':'med'}); le.append(he)
    if dash:
        d=le.makeelement(qn('a:prstDash'),{'val':'dash'}); le.insert(0,d)
    return cn

def node(s,x,y,w,h,t1,t2=None,fill=WHITE,line=LINE2,tc=INK,t2c=MUTE,s1=13,s2=10.5):
    box(s,x,y,w,h,fill,line)
    rows=[[(t1,s1,tc,True,MONO)]]
    if t2: rows.append([(t2,s2,t2c,False,MONO)])
    txt(s,x,y+ (h/2 - (0.32 if t2 else 0.13)),w,h,rows,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.TOP,sp=3)

# ── 1 TITLE ──────────────────────────────────────────────
s=slide(INK2)
box(s,0.42,0.42,W-0.84,H-0.84,None,RGBColor(0x4a,0x58,0x62))  # frame override
txt(s,Mx+0.4,0.6,CW,0.4,[[("SNOWFLAKE DCM · PROOF OF CONCEPT",11,RGBColor(0xB9,0xC4,0xCC),False,MONO)]])
box(s,Mx,0.72,0.28,0.028,AMBLINE)
title(s,[("The database that",WHITE)],y=2.1,size=48)
title(s,[("couldn't see itself",WHITE)],y=3.0,size=48)
box(s,Mx,4.05,0.75,0.04,AMBER)
txt(s,Mx,4.35,7.6,1.2,[
 [("Our repo could always ",18,RGBColor(0xC4,0xCE,0xD5),False,SANS),("build",18,WHITE,True,SANS),(" the database.",18,RGBColor(0xC4,0xCE,0xD5),False,SANS)],
 [("It could never prove the database still ",18,RGBColor(0xC4,0xCE,0xD5),False,SANS),("matched",18,WHITE,True,SANS),(".",18,RGBColor(0xC4,0xCE,0xD5),False,SANS)],
],leading=1.3,sp=2)
txt(s,Mx,6.05,CW,0.4,[[("8",13,AMBLINE,True,MONO),(" tables   ·   ",13,RGBColor(0x94,0xA3,0xAD),False,MONO),("53",13,AMBLINE,True,MONO),(" columns   ·   ",13,RGBColor(0x94,0xA3,0xAD),False,MONO),("13",13,AMBLINE,True,MONO),(" findings   ·   personal account · Aug 2026",13,RGBColor(0x94,0xA3,0xAD),False,MONO)]])
foot(s,"schema drift, caught",1,dark=True)

# ── 2 PROBLEM ────────────────────────────────────────────
s=slide(); kicker(s,"the problem")
title(s,[("Our pipelines cannot ",INK),("see",AMBER),(" the database",INK)])
txt(s,Mx,2.5,CW,0.7,[[("70 ",15,INK,True,SANS),("CREATE TABLE IF NOT EXISTS",13.5,INK2,False,MONO),(" statements across 8 pipeline files. Against a table that already exists, that statement does nothing — and cannot look inside.",15,MUTE,False,SANS)]],leading=1.4)

def stepbox(x, y, w, h, header, lines):
    """titled box: header at top, mono lines stacked below — no overlap."""
    box(s,x,y,w,h,WHITE,LINE2)
    txt(s,x+0.22,y+0.16,w-0.4,0.32,[[(header,13,INK,True,MONO)]])
    box(s,x+0.22,y+0.52,w-0.44,0.014,LINE)          # divider under header
    for i,(t,c,b) in enumerate(lines):
        txt(s,x+0.22,y+0.66+i*0.28,w-0.4,0.3,[[(t,11.5,c,b,MONO)]])

ry=3.55; bw=3.15; gap=0.55; bh=1.85
stepbox(Mx, ry, bw, bh, "1 · created", [
    ("ID    VARCHAR(36)",  MUTE, False),
    ("NAME  VARCHAR(200)", MUTE, False)])
stepbox(Mx+bw+gap, ry, bw, bh, "2 · altered by hand", [
    ("ID    VARCHAR(36)",  MUTE, False),
    ("NAME  VARCHAR(200)", MUTE, False),
    ("SALARY VARCHAR(100)",AMBER, True)])
stepbox(Mx+2*(bw+gap), ry, bw, bh, "3 · re-run DDL", [
    ('"already exists,',   CLEAN, False),
    (' statement succeeded"',CLEAN, False),
    ("SALARY still there", AMBER, True)])
arrow(s,Mx+bw,ry+bh/2,Mx+bw+gap,ry+bh/2)
arrow(s,Mx+2*bw+gap,ry+bh/2,Mx+2*(bw+gap),ry+bh/2)

txt(s,Mx,ry+bh+0.25,CW,0.8,[[('Snowflake is honest — it says it did nothing. The defect is that ',13,MUTE,False,SANS),("“did nothing” counts as success",13,INK,True,SANS),(", and nothing compares the table to the declaration. It stays invisible because it is harmless.",13,MUTE,False,SANS)]],leading=1.4)
foot(s,"the problem",2)

# ── 3 GUARANTEES ─────────────────────────────────────────
s=slide(); kicker(s,"two guarantees")
title(s,[("Only one of these was ever real",INK)])
txt(s,Mx,2.45,CW,0.4,[[('“repo” = the Matillion pipeline repo that holds the 70 DDL statements.',12.5,MUTE,False,MONO)]])
ty=3.1; cw=[7.0,2.0,2.0]
box(s,Mx,ty,sum(cw),0.55,INK2)
for j,htext in enumerate(["","BEFORE","WITH DCM"]):
    txt(s,Mx+sum(cw[:j])+0.2,ty+0.13,cw[j],0.35,[[(htext,11,WHITE,True,MONO)]])
rows=[("The repo can build the database","YES",CLEAN,"YES",CLEAN),
      ("The database still matches the repo","NO",ERROR,"YES",CLEAN)]
for i,(lab,b,bc,d,dc) in enumerate(rows):
    yy=ty+0.55+i*0.62
    if i%2: box(s,Mx,yy,sum(cw),0.62,RGBColor(0xF5,0xF6,0xF8))
    box(s,Mx,yy+0.62-0.012,sum(cw),0.012,LINE)
    txt(s,Mx+0.2,yy+0.16,cw[0],0.4,[[(lab,15,INK,False,SANS)]])
    txt(s,Mx+cw[0]+0.2,yy+0.16,cw[1],0.4,[[(b,15,bc,True,SANS)]])
    txt(s,Mx+cw[0]+cw[1]+0.2,yy+0.16,cw[2],0.4,[[(d,15,dc,True,SANS)]])
txt(s,Mx,ty+2.2,CW,1.3,[[("Two guarantees — only the first was ever true, and nothing could tell you the difference. Why we went looking: a dashboard dimension sat frozen for ",13.5,MUTE,False,SANS),("seven months",13.5,INK,True,SANS),(". A different defect — a trailing comma an alert-then-succeed pattern hid — but the ",13.5,MUTE,False,SANS),("same shape",13.5,AMBER,True,SANS),(": something that looked like a check and never checked.",13.5,MUTE,False,SANS)]],leading=1.45)
foot(s,"two guarantees",3)

# ── 4 WHAT DCM IS ────────────────────────────────────────
s=slide(); kicker(s,"what dcm projects is")
title(s,[("Declare the state.",INK)],y=1.5); title(s,[("Snowflake finds the ",INK),("difference",AMBER),(".",INK)],y=2.35)
by=3.5
for i,(tag,desc) in enumerate([("DEFINE","a description of what should be true — not an instruction"),
                                ("PLAN","a dry run. changes nothing"),
                                ("DEPLOY","applies the changeset. a person's decision")]):
    yy=by+i*0.62
    b=box(s,Mx,yy,1.3,0.42,PANEL); txt(s,Mx,yy+0.1,1.3,0.3,[[(tag,11,INK2,True,MONO)]],align=PP_ALIGN.CENTER)
    txt(s,Mx+1.55,yy+0.07,4.6,0.5,[[(desc,13.5,MUTE,False,SANS)]],leading=1.2)
txt(s,Mx,by+2.05,5.8,0.8,[[("Native to Snowflake · no external state file · GA since Aug 2026 · broad object coverage — ",12,FAINT,False,SANS),("TABLE · VIEW · TASK · ROLE",11,FAINT,False,MONO),(" and more.",12,FAINT,False,SANS)]],leading=1.35)
codebox(s,Mx+6.3,3.35,CW-6.3,[
 [("DEFINE TABLE",INK,True),(" DEVELOP.PRE.DIM_PBI_CAPACITIES (",INK2,False)],
 [("    ID              VARCHAR(36)  ",INK2,False),("NOT NULL",MUTE,False),(",",INK2,False)],
 [("    NAME            VARCHAR(500),",INK2,False)],
 [("    SKU             VARCHAR(50),",INK2,False)],
 [("    IS_CURRENT_FLAG NUMBER(1,0)  ",INK2,False),("DEFAULT 1",MUTE,False)],
 [(");",INK2,False)],
 [("",INK2,False)],
 [("-- remove this block, and the next",MUTE,False)],
 [("-- deploy drops the table.",AMBER,True)],
],size=12.5)
foot(s,"what dcm is",4)

# ── 5 BEFORE/AFTER ───────────────────────────────────────
s=slide(); kicker(s,"infrastructure as code")
title(s,[("The DDL stops living inside the ETL",INK)])
cy=2.8; cw2=(CW-0.5)/2
box(s,Mx,cy,cw2,2.6,WHITE,LINE)
txt(s,Mx+0.3,cy+0.25,cw2-0.6,0.3,[[("BEFORE",11,FAINT,True,MONO)]])
bullets(s,Mx+0.3,cy+0.85,cw2-0.6,[("DDL embedded in orchestration — 70 statements, 8 files",MUTE),("runs as a side effect of a data load",MUTE),("no review · no diff",MUTE),("drift is invisible",ERROR)],size=13,gap=0.42,dot=MUTE)
box(s,Mx+cw2+0.5,cy,cw2,2.6,WHITE,AMBLINE,laccent=AMBER)
txt(s,Mx+cw2+0.8,cy+0.25,cw2-0.6,0.3,[[("AFTER",11,AMBER,True,MONO)]])
bullets(s,Mx+cw2+0.8,cy+0.85,cw2-0.6,[("DDL in git, reviewed like application code",MUTE),("the plan shows impact before it lands",MUTE),("one command rebuilds an environment",MUTE),("drift becomes detectable",CLEAN)],size=13,gap=0.42,dot=AMBER)
txt(s,Mx,cy+2.95,CW,0.6,[[("The database becomes a ",18,INK2,False,SANS),("build output",18,AMBER,True,SANS),(" — not a place things quietly accumulate.",18,INK2,False,SANS)]])
foot(s,"infrastructure as code",5)

# ── 6 FOLDER + MANIFEST ──────────────────────────────────
s=slide(); kicker(s,"project structure")
title(s,[("The layout is ",INK),("fixed",AMBER),(", not a convention",INK)])
codebox(s,Mx,2.75,5.6,[
 [("manifest.yml",INK,True),("        which account · what varies",MUTE,False)],
 [("sources/",INK,True)],
 [("  definitions/",INK,True),("      every object lives here",AMBER,True)],
 [("  macros/",INK2,False),("           optional Jinja",MUTE,False)],
 [("out/",INK2,False),("               generated artifacts",MUTE,False)],
],size=13)
txt(s,Mx,5.0,5.6,0.7,[[("sources/definitions/",11,INK2,False,MONO),(" is required — files anywhere else are ",13,MUTE,False,SANS),("not read",13,INK,True,SANS),(".",13,MUTE,False,SANS)]],leading=1.35)
txt(s,Mx+6.1,2.7,CW-6.1,0.3,[[("THE MANIFEST — WHERE, AND WHAT VARIES",10.5,FAINT,True,MONO)]])
codebox(s,Mx+6.1,3.05,CW-6.1,[
 [("targets:",INK2,False)],
 [("  DCM_DEV:",INK2,False)],
 [("    account_identifier: ORG-ACCOUNT",INK2,False)],
 [("    project_name: …PBI_CAPACITIES",INK2,False)],
 [("templating:",INK2,False)],
 [("  configurations:",INK2,False)],
 [('    DEV:  { env_suffix: "" }',INK2,False)],
 [('    PROD: { env_suffix: ',INK2,False),('"_PROD"',AMBER,True),(' }',INK2,False)],
],size=12)
txt(s,Mx+6.1,5.55,CW-6.1,0.6,[[("Promotion is pointing at a different target. ",13,MUTE,False,SANS),("You copy nothing.",13,INK,True,SANS)]],leading=1.3)
foot(s,"project structure",6)

# ── 7 ARCHITECTURE ───────────────────────────────────────
s=slide(); kicker(s,"architecture")
title(s,[("How it fits together",INK)])
ay=2.7
node(s,Mx,ay,1.9,0.95,"GitHub","authoritative",fill=INK2,tc=WHITE,t2c=RGBColor(0xB9,0xC4,0xCC),s1=13)
node(s,Mx+2.6,ay,2.15,0.95,"GIT REPOSITORY","account clone",s1=12.5)
node(s,Mx+5.4,ay,2.0,0.95,"DCM PROJECT","history · artifacts",s1=12.5)
node(s,Mx+8.15,ay,2.5,0.95,"PROD database","…_PROD · from git",fill=INK2,tc=WHITE,t2c=RGBColor(0xB9,0xC4,0xCC),s1=13)
arrow(s,Mx+1.9,ay+0.47,Mx+2.58,ay+0.47); txt(s,Mx+1.85,ay+0.06,0.9,0.3,[[("FETCH",8.5,MUTE,False,MONO)]],align=PP_ALIGN.CENTER)
arrow(s,Mx+4.75,ay+0.47,Mx+5.38,ay+0.47)
arrow(s,Mx+7.4,ay+0.32,Mx+8.13,ay+0.32); txt(s,Mx+7.25,ay-0.1,1.5,0.3,[[("DEPLOY · from git",8,MUTE,False,MONO)]],align=PP_ALIGN.CENTER)
arrow(s,Mx+7.4,ay+0.64,Mx+8.13,ay+0.64,color=AMBER,dash=True); txt(s,Mx+5.9,ay+0.72,1.9,0.3,[[("PLAN · reads only",8,AMBER,False,MONO)]],align=PP_ALIGN.CENTER)
# nightly monitor row
ly=ay+2.05
node(s,Mx+5.4,ly,2.0,0.9,"Nightly task","05:00 · PLAN only",fill=AMBSOFT,line=AMBLINE,tc=AMBD,t2c=AMBD,s1=12.5)
node(s,Mx+7.7,ly,1.7,0.9,"Drift log","durable record",s1=12.5)
node(s,Mx+9.65,ly,1.4,0.9,"Alert","if not CLEAN",fill=INK2,tc=WHITE,t2c=RGBColor(0xB9,0xC4,0xCC),s1=12.5)
arrow(s,Mx+6.4,ay+0.95,Mx+6.4,ly-0.02); txt(s,Mx+6.55,ly-0.72,1.9,0.3,[[("FETCH → PLAN vs git",8.5,MUTE,False,MONO)]])
arrow(s,Mx+7.4,ly+0.45,Mx+7.68,ly+0.45)
arrow(s,Mx+9.4,ly+0.45,Mx+9.63,ly+0.45,color=AMBER)
txt(s,Mx,H-1.55,CW,0.8,[[("Dev deploys from the workspace (next slide); ",13,MUTE,False,SANS),("prod is deployed from git",13,INK,True,SANS),(", and the nightly check watches prod. PLAN is automated; DEPLOY is a human decision — it drops columns to revert.",13,MUTE,False,SANS)]],leading=1.4)
foot(s,"architecture",7)

# ── 8 THREE COPIES ───────────────────────────────────────
s=slide(); kicker(s,"the distinction people miss")
node(s,Mx+1.0,2.7,2.4,0.9,"GitHub","authoritative",fill=INK2,tc=WHITE,t2c=RGBColor(0xB9,0xC4,0xCC))
node(s,Mx,4.5,2.4,1.05,"GIT REPOSITORY","read by the task")
node(s,Mx+2.9,4.5,2.4,1.05,"Workspace","invisible to automation",t2c=AMBER)
arrow(s,Mx+1.7,3.6,Mx+1.0,4.48); txt(s,Mx+0.7,3.95,0.9,0.3,[[("FETCH",9.5,MUTE,False,MONO)]])
arrow(s,Mx+2.6,3.6,Mx+3.5,4.48); txt(s,Mx+3.5,3.95,0.9,0.3,[[("Pull",9.5,MUTE,False,MONO)]])
txt(s,Mx+0.2,5.75,5.2,0.4,[[("siblings — not parent & child",10.5,FAINT,False,MONO)]],align=PP_ALIGN.CENTER)
tx=Mx+6.0
title(s,[("Three copies.",INK)],y=2.55,size=30,x=tx,w=CW-6.0); title(s,[("All can be ",INK),("stale",AMBER),(".",INK)],y=3.3,size=30,x=tx,w=CW-6.0)
bullets(s,tx,4.25,CW-6.0,[
 ("A GIT REPOSITORY is a snapshot, not a live link — FETCH or it is stale",MUTE),
 ("A workspace lives in your personal database; a scheduled task cannot read it",MUTE),
 ("GitHub is authoritative. The other two are caches.",INK),
],size=14,gap=0.72,dot=AMBER)
foot(s,"three copies",8)

# ── 9 ONE DIRECTION ──────────────────────────────────────
s=slide(); kicker(s,"why drift detection must exist")
title(s,[("A change travels ",INK),("one direction only",AMBER)])
chain=["edit","commit","GitHub","FETCH","PLAN","review","DEPLOY","database"]
n=len(chain); bw=1.28; gap=0.18; startx=Mx+ (CW-(n*bw+(n-1)*gap))/2; cy=3.4
for i,c in enumerate(chain):
    x=startx+i*(bw+gap); hot=c in("GitHub","database")
    node(s,x,cy,bw,0.72,c,None,fill=(INK2 if hot else WHITE),tc=(WHITE if hot else INK),s1=12)
    if i<n-1: arrow(s,x+bw,cy+0.36,x+bw+gap,cy+0.36)
# broken return
rc=s.shapes.add_connector(2,Inches(startx+7*(bw+gap)+bw/2),Inches(cy+0.72),Inches(startx+2*(bw+gap)+bw/2),Inches(cy+0.72))
rc.line.color.rgb=AMBER; rc.line.width=Pt(1.4); rc.shadow.inherit=False
d=rc.line._get_or_add_ln().makeelement(qn('a:prstDash'),{'val':'dash'}); rc.line._get_or_add_ln().insert(0,d)
midx=startx+4.5*(bw+gap)
box(s,midx,cy+1.05,0.3,0.02,AMBER)  # placeholder
txt(s,startx,cy+1.15,n*bw+(n-1)*gap,0.4,[[("✕  no path back",13,AMBER,True,SANS)]],align=PP_ALIGN.CENTER)
txt(s,Mx,H-1.5,CW,0.6,[[("A hand-made ",13.5,MUTE,False,SANS),("ALTER",12,INK2,False,MONO),(" never flows back into the repo. That asymmetry is exactly why the third question has to be asked at all.",13.5,MUTE,False,SANS)]],leading=1.35)
foot(s,"one direction",9)

# ── 10 PROMOTION: DEV FAST, PROD FROM GIT ────────────────
s=slide(); kicker(s,"how a developer ships")
title(s,[("Dev fast. ",INK),("Prod from git.",AMBER)])
node(s,Mx,3.65,2.2,0.85,"workspace","where I edit")
node(s,Mx+4.0,2.7,2.5,0.8,"DEMO_PBI_DEV","fast · no gate",fill=INK2,tc=WHITE,t2c=RGBColor(0xB9,0xC4,0xCC),s1=12.5)
node(s,Mx+4.0,4.55,2.5,0.8,"git · main","commit · PR · review",s1=12.5)
node(s,Mx+8.3,4.55,3.0,0.8,"DEMO_PBI_PROD","reviewed · reproducible",fill=AMBSOFT,line=AMBLINE,tc=AMBD,t2c=AMBD,s1=12.5)
arrow(s,Mx+2.2,3.85,Mx+3.98,3.15); txt(s,Mx+2.35,3.3,1.8,0.3,[[("Deploy · DEV",9.5,MUTE,False,MONO)]])
arrow(s,Mx+2.2,4.1,Mx+3.98,4.8); txt(s,Mx+2.45,4.5,1.0,0.3,[[("push",9.5,MUTE,False,MONO)]])
arrow(s,Mx+6.5,4.95,Mx+8.28,4.95,color=AMBER); txt(s,Mx+6.4,4.62,3.6,0.3,[[("DEPLOY USING CONFIGURATION PROD",9,AMBER,False,MONO)]])
txt(s,Mx,H-1.5,CW,0.8,[[("Same files build both — ",13,MUTE,False,SANS),("env_suffix",11,INK2,False,MONO),(" is the only difference. Prod is deployed ",13,MUTE,False,SANS),("from git, never the workspace",13,INK,True,SANS),(" — so git and prod agree by construction, and the drift check stays honest.",13,MUTE,False,SANS)]],leading=1.4)
foot(s,"promotion",10)

# ── 11 THREE DIFFS ───────────────────────────────────────
s=slide(); kicker(s,"the spine of the whole idea")
title(s,[("Three questions. Git answers two.",INK)])
qy=2.9
data=[("01","What changed in the definition?","git diff",False),
      ("02","What will change in the database?","DCM PLAN",False),
      ("03","What changed without going through either?","drift check",True)]
for i,(nn,q,tool,hot) in enumerate(data):
    yy=qy+i*0.92
    box(s,Mx,yy,CW,0.78,WHITE,(AMBLINE if hot else LINE),laccent=(AMBER if hot else None))
    txt(s,Mx+0.35,yy+0.2,1.0,0.4,[[(nn,18,AMBER,True,MONO)]])
    txt(s,Mx+1.5,yy+0.19,CW-4.5,0.5,[[(q,17,INK,(True if hot else False),SANS)]],anchor=MSO_ANCHOR.MIDDLE)
    tb=box(s,W-Mx-2.3,yy+0.22,1.9,0.34,(DRIFTS if hot else PANEL))
    txt(s,W-Mx-2.3,yy+0.3,1.9,0.3,[[(tool,11.5,(DRIFT if hot else INK2),True,MONO)]],align=PP_ALIGN.CENTER)
txt(s,Mx,qy+3.0,CW,0.6,[[("The third is the one we have ",18,INK2,False,SANS),("never been able to see.",18,AMBER,True,SANS)]])
foot(s,"three diffs",11)

# ── 11 VERDICTS ──────────────────────────────────────────
s=slide(); kicker(s,"drift detection")
# tree left
node(s,Mx+1.0,2.7,2.4,0.6,"PLAN compiles?",None,s1=12.5)
node(s,Mx+3.8,4.0,1.9,0.75,"ERROR","un-revertible",fill=ERRORS,line=ERROR,tc=ERROR,t2c=ERROR,s1=13,s2=9.5)
node(s,Mx+0.8,4.0,2.3,0.6,"changeset empty?",None,s1=11.5)
node(s,Mx,5.35,1.9,0.75,"CLEAN","matches repo",fill=CLEANS,line=CLEAN,tc=CLEAN,t2c=CLEAN,s1=13,s2=9.5)
node(s,Mx+2.3,5.35,1.9,0.75,"DRIFT","names the column",fill=DRIFTS,line=DRIFT,tc=DRIFT,t2c=DRIFT,s1=13,s2=9.5)
arrow(s,Mx+3.4,3.0,Mx+4.5,3.98); txt(s,Mx+3.85,3.35,0.6,0.3,[[("no",9.5,ERROR,False,MONO)]])
arrow(s,Mx+2.0,3.3,Mx+1.9,3.98); txt(s,Mx+2.05,3.5,0.5,0.3,[[("yes",9.5,MUTE,False,MONO)]])
arrow(s,Mx+1.3,4.6,Mx+0.95,5.33); txt(s,Mx+0.55,4.9,0.7,0.3,[[("empty",9,CLEAN,False,MONO)]])
arrow(s,Mx+2.4,4.6,Mx+3.0,5.33); txt(s,Mx+2.85,4.9,0.6,0.3,[[("not",9,DRIFT,False,MONO)]])
# right text
tx=Mx+6.2
title(s,[("Three verdicts,",INK)],y=2.6,size=30,x=tx,w=CW-6.2); title(s,[("not two",INK)],y=3.35,size=30,x=tx,w=CW-6.2)
txt(s,tx,4.25,CW-6.2,0.6,[[("A two-state monitor — “changeset empty or not” — reads the worst case as a broken job.",14,MUTE,False,SANS)]],leading=1.35)
box(s,tx,5.05,CW-6.2,1.4,WHITE,LINE); box(s,tx,5.05,0.045,1.4,AMBER)
txt(s,tx+0.32,5.24,1.3,0.6,[[("45",30,AMBER,True,MONO)]])
txt(s,tx+1.5,5.34,CW-7.7,0.5,[[("of 53",13,FAINT,False,SANS),(" columns are not the",13,INK2,False,SANS)],[("last column in their table.",13,INK2,False,SANS)]],sp=1,leading=1.15)
txt(s,tx+0.32,6.0,CW-6.6,0.4,[[("Restoring one mid-list is a reorder Snowflake can't do — ERROR is the common case.",12,MUTE,False,SANS)]],leading=1.2)
foot(s,"verdicts",12)

# ── 12 RETURNS + AUDIT ───────────────────────────────────
s=slide(); kicker(s,"what it returns · what snowflake forgets")
title(s,[("It names the column",INK)],y=1.5,size=27)
codebox(s,Mx,2.35,5.7,[
 [("ALTER TABLE",INK,True),(" …PRE.DIM_PBI_CAPACITIES",INK2,False)],
 [("  columns:",INK2,False)],
 [('    removed "HAND_ADDED_BY_A_HUMAN"',AMBER,True)],
 [("            VARCHAR(100), nullable",AMBER,True)],
],size=12.5)
txt(s,Mx,4.15,5.7,1.0,[[("Not “a table differs.” The name, the type, and which way the fix runs — the difference between an alert worth waking for and one people mute.",13,MUTE,False,SANS)]],leading=1.4)
txt(s,Mx+6.1,1.55,CW-6.1,0.3,[[("MEASURED — 7+ PLANS, 4 DEPLOYS",10.5,FAINT,True,MONO)]])
codebox(s,Mx+6.1,1.9,CW-6.1,[
 [("SELECT PHASE, COUNT(*)",INK2,False)],
 [("  FROM DCM_DEPLOYMENT_HISTORY(…)",INK2,False)],
 [("",INK2,False)],
 [("  PHASE     N",MUTE,False)],
 [("  DEPLOY    4",INK,True),("   ← no PLAN rows",AMBER,True)],
],size=12.5)
txt(s,Mx+6.1,3.9,CW-6.1,1.4,[[("The drift check is the one thing Snowflake doesn't record.",13,INK,True,SANS),(" Deployments keep 12-month artifacts; plans keep nothing, with no ",13,MUTE,False,SANS),("ACCOUNT_USAGE",11,INK2,False,MONO),(" view — so we keep our own log. It answers ",13,MUTE,False,SANS),("when did this drift start?",13,INK2,True,SANS)]],leading=1.4)
foot(s,"the record",13)

# ── 13 FINDINGS ──────────────────────────────────────────
s=slide(); kicker(s,"what we found by breaking it")
title(s,[("Thirteen findings. Four from deliberate sabotage.",INK)])
fy=2.8
rows=[("F5","Drift is detected and reported at column level — the verdict"),
      ("F6","DCM records DEPLOY, never PLAN — an audit table is mandatory"),
      ("F7","Dropping a column that is not last is un-revertible"),
      ("F9","Our own monitor sat suspended while its health view reported OK"),
      ("F11","Has now run unattended two consecutive nights, ~30s per run")]
for i,(fid,desc) in enumerate(rows):
    yy=fy+i*0.58
    if i%2: box(s,Mx,yy,CW,0.58,RGBColor(0xF5,0xF6,0xF8))
    box(s,Mx,yy+0.58-0.01,CW,0.01,LINE)
    txt(s,Mx+0.2,yy+0.15,1.0,0.3,[[(fid,14,AMBER,True,MONO)]])
    txt(s,Mx+1.2,yy+0.15,CW-1.4,0.35,[[(desc,14.5,INK,False,SANS)]])
txt(s,Mx,fy+3.2,CW,0.6,[[("Every failure we found was something ",17,INK2,False,SANS),("reporting success without checking",17,AMBER,True,SANS),(" what it claimed to check.",17,INK2,False,SANS)]])
foot(s,"findings",14)

# ── 14 DEMO ──────────────────────────────────────────────
s=slide(); kicker(s,"live · end to end · ~20 min")
title(s,[("What we'll demo",INK)])
items=[(1,"CREATE TABLE IF NOT EXISTS succeeding against a table it never looked at",False),
       (2,"Build a project from nothing — scaffold, definitions, plan, deploy",False),
       (3,"Push to git — the repo fills up",False),
       (4,"Snowflake pulls its own copy; the plan agrees from both directions",False),
       (5,"Add a table, push, fetch, deploy — a change travelling through git",False),
       (6,"Tamper by hand — GitHub shows nothing, the drift check names it",True),
       (7,"The un-revertible case, and why it changes the design",False)]
colw=(CW-0.6)/2
for i,(nn,desc,hot) in enumerate(items):
    col=i//4; row=i%4
    x=Mx+col*(colw+0.6); yy=2.9+row*0.82
    if hot: box(s,x,yy-0.05,colw,0.72,AMBSOFT)
    txt(s,x+0.1,yy,0.5,0.4,[[(str(nn),15,(AMBER if hot else FAINT),True,MONO)]])
    txt(s,x+0.7,yy,colw-0.8,0.7,[[(desc,13,(INK2 if hot else MUTE),(True if hot else False),SANS)]],leading=1.25)
foot(s,"the demo",15)

# ── 15 BEYOND SCHEMA: CONCEPT ────────────────────────────
s=slide(); kicker(s,"beyond schema · what's next")
node(s,Mx,3.0,2.7,0.85,"does it MATCH?","PLAN · structure",s1=12.5,s2=10)
node(s,Mx,4.7,2.7,0.85,"does the data PASS?","snow dcm test",fill=AMBSOFT,line=AMBLINE,tc=AMBD,t2c=AMBD,s1=12.5,s2=10)
node(s,Mx+4.3,3.85,2.4,0.95,"the table","PRE.DIM_PBI_…",fill=INK2,tc=WHITE,t2c=RGBColor(0xB9,0xC4,0xCC))
arrow(s,Mx+2.7,3.42,Mx+4.28,4.05)
arrow(s,Mx+2.7,5.12,Mx+4.28,4.55,color=AMBER)
txt(s,Mx+2.9,3.95,1.5,0.3,[[("this POC",9.5,FAINT,False,MONO)]])
txt(s,Mx+2.9,4.62,1.7,0.3,[[("not yet used",9.5,AMBER,False,MONO)]])
tx=Mx+7.2
title(s,[("Gate the ",INK),("data",AMBER),(",",INK)],y=2.7,size=28,x=tx,w=CW-7.2)
title(s,[("not just the shape",INK)],y=3.45,size=28,x=tx,w=CW-7.2)
txt(s,tx,4.4,CW-7.2,1.1,[[("This POC proved the database matches the repo. The same declarative project can prove the data passes — quality checks declared beside the schema.",13.5,MUTE,False,SANS)]],leading=1.4)
txt(s,tx,5.75,CW-7.2,1.0,[[("ATTACH a data metric function with an EXPECTATION, then ",13.5,MUTE,False,SANS),("snow dcm test",11.5,INK2,False,MONO),(". A GA capability we never used.",13.5,INK,True,SANS)]],leading=1.4)
foot(s,"beyond schema",16)

# ── 16 BEYOND SCHEMA: EXAMPLES ───────────────────────────
s=slide(); kicker(s,"beyond schema · examples")
title(s,[("Quality gates, declared beside the schema",INK)])
codebox(s,Mx,2.75,6.5,[
 [("-- in the same sources/definitions file",MUTE,False)],
 [("ATTACH DATA METRIC FUNCTION",INK,True),(" SNOWFLAKE.CORE.NULL_COUNT",INK2,False)],
 [("  TO TABLE",INK,True),(" DEVELOP.PRE.DIM_PBI_CAPACITIES ",INK2,False),("ON",INK,True),(" (ID)",INK2,False)],
 [("  EXPECTATION",INK,True),(" no_null_ids ( ",INK2,False),("VALUE = 0",AMBER,True),(" );",INK2,False)],
 [("",INK2,False)],
 [("ATTACH DATA METRIC FUNCTION",INK,True),(" SNOWFLAKE.CORE.DUPLICATE_COUNT",INK2,False)],
 [("  TO TABLE",INK,True),(" DEVELOP.PRE.DIM_PBI_CAPACITIES ",INK2,False),("ON",INK,True),(" (ID)",INK2,False)],
 [("  EXPECTATION",INK,True),(" unique_ids ( ",INK2,False),("VALUE = 0",AMBER,True),(" );",INK2,False)],
],size=11)
txt(s,Mx,5.15,6.5,0.7,[[("The left side is always the keyword ",12.5,MUTE,False,SANS),("VALUE",11,INK2,False,MONO),(". No subqueries, casts or arithmetic — a check, not a query.",12.5,MUTE,False,SANS)]],leading=1.35)
txt(s,Mx+6.95,2.7,CW-6.95,0.3,[[("ONE COMMAND RUNS THEM ALL",10.5,FAINT,True,MONO)]])
codebox(s,Mx+6.95,3.05,CW-6.95,[
 [("$ snow dcm test PBI_CAPACITIES",INK2,False)],
 [("",INK2,False)],
 [("  PASS  no_null_ids",CLEAN,True)],
 [("  FAIL  unique_ids",AMBER,True)],
 [("     Expected VALUE = 0, got 1",MUTE,False)],
 [("  1 passed, 1 failed",MUTE,False)],
],size=12)
txt(s,Mx+6.95,5.25,CW-6.95,1.3,[[("A ",12.5,MUTE,False,SANS),("verified run",12.5,INK,True,SANS),(", not a mock-up. Exits non-zero on failure, so it gates a pipeline. The same check sits on landing, staging and presentation — bad data caught where it enters, not three dashboards later.",12.5,MUTE,False,SANS)]],leading=1.4)
foot(s,"beyond schema",17)

# ── 15 CLOSE ─────────────────────────────────────────────
s=slide(); kicker(s,"honest limits")
title(s,[("What it does",INK)],y=1.6,size=30); title(s,[("not yet prove",INK)],y=2.35,size=30)
bullets(s,Mx,3.5,5.6,[
 ("One slice — 8 tables, 53 columns, no data",MUTE),
 ("Untested at full-estate scale (52 tables)",MUTE),
 ("Data-quality gates (expectations) not yet used",MUTE),
 ("Grants, tasks, streams not yet tested by us",MUTE),
 ("ERROR recovery rehearsed only on empty tables",MUTE)],size=14,gap=0.52,dot=MUTE)
box(s,Mx+6.4,3.3,0.012,2.8,LINE)
txt(s,Mx+6.8,3.5,CW-6.8,2.0,[
 [("The database can now tell us ",18,INK2,False,SANS),("every morning",18,AMBER,True,SANS),(" whether it still matches the repo —",18,INK2,False,SANS)],
 [("and name the column when it doesn't.",18,INK2,False,SANS)],
],leading=1.3,sp=10)
txt(s,Mx+6.8,5.7,CW-6.8,0.6,[[("Scheduled PLAN first — it's read-only. DEPLOY stays a human decision.",12.5,FAINT,False,SANS)]],leading=1.3)
foot(s,"schema drift, caught",18)

prs.save("DCM_Presentation.pptx")
print("saved", len(prs.slides._sldIdLst), "slides")

# ── layout self-check (run after save) ───────────────────────────────
def _validate(path="DCM_Presentation.pptx"):
    from pptx import Presentation as _P
    from pptx.enum.shapes import MSO_SHAPE_TYPE as _T
    import math as _m
    p=_P(path); EMU=914400; SW=p.slide_width; SH=p.slide_height; issues=0
    def er(sh):
        L=sh.left/EMU;Tp=sh.top/EMU;W=(sh.width or 0)/EMU;h=0;mw=0
        for pa in sh.text_frame.paragraphs:
            rs=pa.runs
            if not rs: continue
            t="".join(r.text for r in rs); f=max([(r.font.size.pt if r.font.size else 12) for r in rs]+[12])
            cpl=max(1,W*72/(0.52*f)); ln=max(1,_m.ceil(len(t)/cpl)); h+=ln*f*1.32/72; mw=max(mw,min(W,len(t)*0.52*f/72))
        return (L,Tp,L+mw,Tp+h)
    def ar(r):return max(0,r[2]-r[0])*max(0,r[3]-r[1])
    def it(a,b):return ar((max(a[0],b[0]),max(a[1],b[1]),min(a[2],b[2]),min(a[3],b[3])))
    for i,s in enumerate(p.slides,1):
        tb=[sh for sh in s.shapes if sh.has_text_frame and sh.shape_type==_T.TEXT_BOX and sh.text_frame.text.strip()]
        R=[er(sh) for sh in tb]
        for a in range(len(R)):
            for b in range(a+1,len(R)):
                o=it(R[a],R[b]); sm=min(ar(R[a]),ar(R[b]))
                if sm>0.02 and o/sm>0.30: issues+=1; print(f"  OVERLAP slide {i}: {tb[a].text_frame.text[:20]!r} x {tb[b].text_frame.text[:20]!r}")
        FR=0.42*EMU
        for sh in s.shapes:
            if sh.left is None: continue
            R=sh.left+(sh.width or 0); B=sh.top+(sh.height or 0)
            if sh.left<-9000 or sh.top<-9000 or R>SW+9000 or B>SH+9000:
                issues+=1; print(f"  OFF-CANVAS slide {i}")
            elif sh.left<FR-4000 or sh.top<FR-4000 or R>SW-FR+4000 or B>SH-FR+4000:
                t=sh.text_frame.text[:18] if sh.has_text_frame else "shape"
                issues+=1; print(f"  OUT-OF-FRAME slide {i}: {t!r}")
    print("layout check:", "CLEAN" if issues==0 else f"{issues} ISSUES")
    return issues

if __name__=="__main__":
    _validate()
